import os
import dotenv
import json
import copy
import nest_asyncio
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from io import BytesIO
from autogen import UserProxyAgent, AssistantAgent
from autogen.agentchat.contrib.llamaindex_conversable_agent import LLamaIndexConversableAgent
from llama_parse import LlamaParse
from llama_index.core import Settings
from llama_index.core.agent import ReActAgent
from llama_index.core.schema import TextNode
from llama_index.core import VectorStoreIndex
from llama_index.core.node_parser import MarkdownElementNodeParser
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI
from llama_index.postprocessor.flag_embedding_reranker import FlagEmbeddingReranker

nest_asyncio.apply()

# Load environment variables
dotenv.load_dotenv()

# Load API key from environment
OPENAI_API_KEY = os.getenv('TERTIARY_INFOTECH_API_KEY')

# AutoGen Configs
config_list = [
    {
        "model": "gpt-4o-mini",
        "api_key": OPENAI_API_KEY
    }
]

llm_config = {
    "config_list": config_list,
    "timeout": 120,
    "temperature": 0
}

# HARDCODED CONTEXT TO BE REPLACED WITH INTERPRETER AGENT
context = {
    "Course_Title": "Java Programming Methodologies",
    "TGS_Ref_No": "TGS-2024048313",
    "TSC_Title": "Data Governance",
    "TSC_Code": "ICT-DES-4005-1.1",
    "Abilities": [
        "A1: Create a software design blueprint based on a broad design concept, and business and user requirements",
        "A2: Recommend appropriate standards, methods and tools for the design of software, in line with the organisation's practice and design principles",
        "A3: Design functional specifications of software systems to address business and user needs",
        "A4: Evaluate trade offs from the incorporation of different elements into the design, and their impact on overall functionality, interoperability, efficiency and costs of the software",
        "A5: Produce design documentation for complex software",
        "A6: Review design documentations produced"
    ],
    "Knowledge": [
        "K1: Components and requirements of a software design blueprint",
        "K2: Software design standards, methods and tools - and their pros, cons and applications",
        "K3: Requirements of functional specifications of software",
        "K4: Impact of different software design elements on overall software operations and usability"
    ],
    "Learning_Outcomes": [
        "LO1: Develop a software design blueprint based on Java programming integrating design concepts with user requirements.",
        "LO2: Develop software design standards using Java Object Oriented Programming methodologies.",
        "LO3: Design Java data structures and API for software systems tailored to meet business and user needs.",
        "LO4: Evaluate the issues involved in Java applications using debugging and exception handling tools.",
        "LO5: Produce design documentation using Java generics to ensure conformity to technical standards."
    ],
    "Course_Outline": {
        "Topic 1: Introduction to Java": [
            "Editors and Tools",
            "Basic Syntax",
            "Language Syntax Properties",
            "Variables & Datatypes & Literals",
            "Operators",
            "Autoboxing and Unboxing",
            "Enums",
            "Arrays",
            "Strings",
            "Date and Time"
        ],
        "Topic 2: Control Flow": [
            "Statements, Expressions & Blocks",
            "Flow Control statements",
            "Ternary Operator",
            "Loops statements",
            "Nested Loops statements",
            "Loop Control Statements"
        ],
        "Topic 3: Object Oriented Programming": [
            "Scope",
            "Classes & Object",
            "Methods",
            "Constructors",
            "Access Modifiers",
            "‘this’ keyword",
            "Passing by Value",
            "Encapsulation",
            "Inheritance",
            "Abstraction",
            "Interface",
            "Polymorphism"
        ],
        "Topic 4: Data Structures": [
            "Static & Dynamic Array",
            "N-Dimensional Array",
            "Basic Operations on Arrays",
            "Basic operations on Linked List",
            "Arrays & Linked List",
            "Types of Linked List",
            "Stacks & Queues"
        ],
        "Topic 5: Developing an API": [
            "Design the API architecture",
            "Develop the API",
            "Test the API",
            "Monitor the API and iterate on feedback"
        ],
        "Topic 6: Debugging Java Applications": [
            "What is Debugging?",
            "Examining the code",
            "Setting breakpoints",
            "Running the program in debug mode",
            "Analyze the program state",
            "Step through the program",
            "Stopping the debugging session and rerun the program"
        ],
        "Topic 7: Exception Handling": [
            "Exception keywords",
            "Nested exceptions",
            "Throwing exceptions",
            "Exception propagation",
            "Throws clause",
            "Custom exceptions",
            "Chaining exceptions",
            "Exceptions with polymorphism"
        ],
        "Topic 8: File Operations": [
            "File paths",
            "File metadata",
            "Creating regular and temporary files",
            "The try-with-resources statement",
            "Checking if a File or Directory exists",
            "File access modes",
            "Reading from and writing to files"
        ],
        "Topic 9: Using Generics": [
            "Generic types",
            "Bounded type parameters",
            "Inheritance and subtypes",
            "Type inference",
            "Wildcards",
            "Restrictions"
        ],
        "Topic 10: Multi-threading": [
            "Life cycle of a thread",
            "Synchronization",
            "Issues with Multi-threading",
            "Interrupting Threads"
        ]
    },
    "Assessment_Methods": ["Written Assessment (SAQ)", "Practical Performance (PP)"]
}

# Create the content generator agent
slides_writer = AssistantAgent(
    name="Slides_Generator",
    system_message=f"""
    You are a PowerPoint slide writer.
    You have received a JSON object as context data: {context}.
    Your task is to pass this JSON object **exactly as received**, as the `context` parameter, to the `generate_slides` tool function.

    **Important Instructions:**
    - Do not modify or unpack the JSON object.
    - Pass the JSON object **directly** into the `generate_slides` function as follows:
      ```python
      generate_slides(context = context dictionary)
      ```
    - Replace `JSON_Context` with the actual JSON object received.
    Upon successful completion of the task, return `TERMINATE`.
    """,
    llm_config=llm_config,
)

# Create the user proxy agent
user_proxy = UserProxyAgent(
    name="User_Proxy",
    human_input_mode="NEVER",
    max_consecutive_auto_reply=2,
    # is_termination_msg=lambda x: x.get("content", "").find("TERMINATE") >= 0,
    code_execution_config={"work_dir": "output", "use_docker": False} 
)

@user_proxy.register_for_execution()
@slides_writer.register_for_llm(name="generate_slides", description="Generate the Master Slides")
def generate_slides(context: dict) -> str:
    print(context)
    # Function to replace placeholders
    def replace_placeholder(text, replacements):
        for placeholder, replacement in replacements.items():
            if isinstance(replacement, list):
                replacement = "\n".join(replacement)
            if isinstance(replacement, dict):
                # Skip dictionary replacements like <COURSE_OUTLINE>
                continue
            text = text.replace(placeholder, replacement)
        return text

    # Function to add bullet points from a list to a text frame
    def add_bullet_points(text_frame, bullet_points):
        text_frame.clear()  # Clear any existing text
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Level 0 for top-level bullet points

    def get_slide_number_placeholder_properties():
        # Positioning and formatting based on the data provided in the images
        position = (Cm(23.54), Cm(12.95), Cm(1.52), Cm(1.09))  # Left, Top, Width, Height in centimeters
        font_properties = {
            'font_size': Pt(10),
            'font_name': "Arial",
            'font_color': RGBColor(89, 89, 89),  # Dark Gray, Background 2
            'alignment': 'middle',  # Middle alignment
            'horizontal_alignment': 'right',  # Right horizontal alignment
            'autofit': False,
            'wrap_text': True,
            'margins': (Cm(0.25), Cm(0.25), Cm(0.25), Cm(0.25))  # Left, Right, Top, Bottom margins in centimeters
        }
        return position, font_properties

    def move_slides(prs1, prs2):
        subtopic_slides = []
        # Step 1: Find the slide in prs1 with the "Subtopic Placeholder"
        placeholder_index = None
        for idx, slide in enumerate(prs1.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and "Topics" in shape.text:
                    placeholder_index = idx
                    break
            if placeholder_index is not None:
                break

        # If placeholder is found, proceed
        if placeholder_index is not None:
            # Step 2: Collect all slides after the placeholder
            slides_after_placeholder = [prs1.slides[i] for i in range(placeholder_index + 1, len(prs1.slides))]

            # Step 3: Remove slides after the placeholder (starting from the end)
            for i in range(len(prs1.slides) - 1, placeholder_index, -1):
                rId = prs1.slides._sldIdLst[i].rId
                prs1.part.drop_rel(rId)
                del prs1.slides._sldIdLst[i]

            # Step 4: Insert each slide from prs2 after the placeholder slide
            for slide in prs2.slides:
                # Use the layout with slide number placeholder
                new_slide = prs1.slides.add_slide(slide.slide_layout)
                for shape in slide.shapes:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # For images, use add_picture with a BytesIO wrapper for the image blob
                        image_stream = BytesIO(shape.image.blob)
                        new_slide.shapes.add_picture(
                            image_stream,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )
                    else:
                        # Copy other shapes as before
                        new_shape = copy.deepcopy(shape.element)
                        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

                # Attempt to set the title text if it exists
                try:
                    new_slide.shapes.title.text = slide.shapes.title.text
                except AttributeError:
                    pass

                # Record the new slide as a subtopic slide
                subtopic_slides.append(new_slide)

            # Step 5: Re-add the slides that were after the placeholder
            for slide in slides_after_placeholder:
                new_slide = prs1.slides.add_slide(slide.slide_layout)
                for shape in slide.shapes:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # Copy image shapes with BytesIO for image data
                        image_stream = BytesIO(shape.image.blob)
                        new_slide.shapes.add_picture(
                            image_stream,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )
                    else:
                        # Copy other shapes
                        new_shape = copy.deepcopy(shape.element)
                        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

                # Set title text if it exists
                try:
                    new_slide.shapes.title.text = slide.shapes.title.text
                except AttributeError:
                    pass

        else:
            # If no placeholder is found, proceed as needed
            subtopic_slides = []

        return prs1, subtopic_slides

    # Modify the update_slide_numbers function to set alignment correctly
    def update_slide_numbers(presentation, subtopic_slides, slide_number_position, font_properties):
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_number_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    slide_number_placeholder = shape
                    break
            if slide_number_placeholder:
                # Update the placeholder text with the correct slide number
                slide_number_placeholder.text = str(idx)
            elif slide in subtopic_slides and slide_number_position is not None:
                # Use the extracted position and formatting to add the slide number
                left, top, width, height = slide_number_position
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = str(idx)

                # Apply the extracted font formatting
                font = p.font
                font.size = font_properties['font_size']
                font.name = font_properties['font_name']
                font.color.rgb = font_properties['font_color']
                tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = font_properties['margins']
                tf.word_wrap = font_properties['wrap_text']
                
                # Set auto_size based on font_properties['autofit']
                tf.auto_size = MSO_AUTO_SIZE.NONE if not font_properties['autofit'] else MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                # Set vertical alignment to middle
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Set paragraph alignment to right
                p.alignment = PP_ALIGN.RIGHT

    def remove_unused_placeholders(presentation):
        for slide in presentation.slides:
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.is_placeholder:
                    # Check if the placeholder is empty
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if not text:
                            shapes_to_remove.append(shape)
                    elif shape.placeholder_format.type == MSO_SHAPE_TYPE.PICTURE:
                        # For picture placeholders, check if an image is inserted
                        if not shape.has_image:
                            shapes_to_remove.append(shape)
                    else:
                        # Other placeholder types can be checked similarly
                        shapes_to_remove.append(shape)
            # Remove the unused placeholders
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
                
    print("Starting generate_slides function")
    print(f"Received context: {context}")

    # Convert all keys in replacements to uppercase and wrap them with < >
    context = {f"<{key.upper()}>": value for key, value in context.items()}
    print(context)
    try:
        # Load JSON data from file
        with open('Slides/generated_content.json', 'r') as f:
            data = json.load(f)

        # Load the PowerPoint template
        template_path = 'Slides/templates/WSQ/(Template) WSQ - Master Trainer Slides - Course Title - version.pptx'
        prs1 = Presentation(template_path)  # Main presentation

        content_slide_layout = prs1.slide_layouts[1]
        title_slide_layout = prs1.slide_masters[0].slide_layouts.get_by_name("SECTION_HEADER")

        # Step 1: Replace placeholders on existing slides in prs1
        for slide in prs1.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "<COURSE_OUTLINE>" in shape.text: 
                        # Clear the placeholder's content
                        slide_index= prs1.slides.index(slide)
                        break  # Stop looking once <COURSE_OUTLINE> placeholder is replaced
                    if "<LEARNING_OUTCOMES>" in shape.text:
                        add_bullet_points(shape.text_frame, context["<LEARNING_OUTCOMES>"])
                    elif "<ASSESSMENT_METHODS>" in shape.text:
                        add_bullet_points(shape.text_frame, context["<ASSESSMENT_METHODS>"])
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.text = replace_placeholder(paragraph.text, context)
        
        # Remove the placeholder slide
        # xml_slides = prs1.slides._sldIdLst
        # slides = list(xml_slides)
        # xml_slides.remove(slides[slide_index])

        # Step 2: Generate Course Outline Slides from <COURSE_OUTLINE>
        course_outline = context["<COURSE_OUTLINE>"]
        topics_per_slide = 2  # Define how many topics to include per slide
        topic_list = list(course_outline.items())  # Convert dictionary to a list of tuples for indexing

        for i in range(0, len(topic_list), topics_per_slide):
            # Add a new slide
            new_slide = prs1.slides.add_slide(content_slide_layout)
            
            # Add a title for the slide
            title_placeholder = new_slide.shapes.title
            title_placeholder.text = "Course Outline"

            # Add topics and their subtopics to the slide
            content_placeholder = new_slide.placeholders[1]  # Placeholder for content
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear any existing content
            
            # Add topics for this slide
            for topic_title, subtopics in topic_list[i:i + topics_per_slide]:
                # Add topic title as a bold heading
                topic_paragraph = text_frame.add_paragraph()
                topic_paragraph.text = topic_title
                topic_paragraph.font.bold = True

                # Add subtopics as bullet points under the topic
                for subtopic in subtopics:
                    subtopic_paragraph = text_frame.add_paragraph()
                    subtopic_paragraph.text = subtopic
                    subtopic_paragraph.level = 1  # Indent bullet points for subtopics
        
        # Step 2: Create a new presentation prs2 for subtopic slides
        prs2 = Presentation()

        # Generate subtopic slides and add them to prs2
        # Generate topic title slides and subtopic slides and add them to prs2
        for topic in data['Topics']:
            # Add a topic title slide with layout 4
            topic_slide = prs2.slides.add_slide(title_slide_layout)
            
            # Set the title for the topic slide
            topic_title_placeholder = topic_slide.shapes.title
            topic_title_placeholder.text = topic['topic']

            for subtopic in topic['subtopics']:
                # Add a slide for each subtopic
                slide = prs2.slides.add_slide(content_slide_layout)
                
                # Set the title for the slide
                title_placeholder = slide.shapes.title
                title_placeholder.text_frame.clear()
                title_placeholder.text = subtopic['subtopic']
                
                # Set up the content placeholder with bullet points
                content_placeholder = slide.placeholders[1]
                content_placeholder.text_frame.clear()  # Clear existing content if any
                
                # Loop through each keypoint and its bullets
                for keypoint in subtopic['keypoints']:
                    if keypoint['keypoint'] == "NO CONTEXT":
                        continue
                    # Add keypoint title as a bold line
                    p = content_placeholder.text_frame.add_paragraph()
                    p.text = keypoint['keypoint']
                    p.font.bold = True
                    
                    # Add each valid bullet point under the keypoint
                    for bullet in keypoint['bullets']:
                        if bullet != "NO CONTEXT":
                            bullet_paragraph = content_placeholder.text_frame.add_paragraph()
                            bullet_paragraph.text = bullet

        # Step 3: Move all slides from prs2 into prs1 after the placeholder slide
        prs1, subtopic_slides = move_slides(prs1, prs2)

        # Step 4: Get the slide number placeholder properties (position and formatting)
        slide_number_position, font_properties = get_slide_number_placeholder_properties()

        # Step 5: Update slide numbers using the extracted properties
        update_slide_numbers(prs1, subtopic_slides, slide_number_position, font_properties)

        # Step 6: Remove unused placeholders in prs1
        remove_unused_placeholders(prs1)

        # Step 7: Save the final presentation with subtopic slides inserted
        prs1.save('Slides/output/autogen_output.pptx')
        print(f'Slides successfully generated to output path: Slides/output/autogen_output.pptx')

    except Exception as e:
        print(f"Error in generate_slides: {str(e)}")
        raise


# Load JSON data from file
with open('Slides/generated_content.json', 'r') as f:
    data = json.load(f)

user_proxy.initiate_chat(
    slides_writer,
    message=f"Please generate the powerpoint slides",
)
