import os
import dotenv
import json
import copy
import nest_asyncio
from docx import Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import Table
from docx.text.paragraph import Paragraph
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from io import BytesIO
from autogen import UserProxyAgent, AssistantAgent
from autogen.agentchat.contrib.llamaindex_conversable_agent import LLamaIndexConversableAgent
from llama_parse import LlamaParse
from llama_index.core import Settings
from llama_index.core.agent import ReActAgent
from llama_index.core.schema import TextNode
from llama_index.core import VectorStoreIndex
from llama_index.core.node_parser import MarkdownElementNodeParser
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI
from llama_index.postprocessor.flag_embedding_reranker import FlagEmbeddingReranker

nest_asyncio.apply()

# Load environment variables
dotenv.load_dotenv()

# Load API key from environment
OPENAI_API_KEY = os.getenv('TERTIARY_INFOTECH_API_KEY')
LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")

# AutoGen Configs
config_list = [
    {
        "model": "gpt-4o-mini",
        "api_key": OPENAI_API_KEY
    }
]

llm_config = {
    "config_list": config_list,
    "timeout": 120,
    "temperature": 0
}

# LlamaIndex Configs
llm = OpenAI(
model="gpt-4o-mini",
temperature=0.0,
api_key=OPENAI_API_KEY,
)

embed_model = OpenAIEmbedding(
model="text-embedding-ada-002",
temperature=0.0,
api_key=OPENAI_API_KEY,
)

def saveToJSON(output_path, data):
    with open(output_path, 'w', encoding='utf-8') as json_file:
        json.dump(data, json_file, ensure_ascii=False, indent=4)
        
# def parse_fg_document(input_dir):
#     """
#     Parse a Facilitator Guide document into a structured JSON format.

#     Args:
#         input_dir (str): Path to the input Word document.

#     Returns:
#         dict: Parsed document content.
#     """
#     # Load the document
#     doc = Document(input_dir)

#     # Initialize containers
#     data = {
#         "Facilitator_Guide": {}
#     }

#     # Function to parse tables with advanced duplication check
#     def parse_table(table):
#         rows = []
#         for row in table.rows:
#             # Process each cell and ensure unique content within the row
#             cells = []
#             for cell in row.cells:
#                 cell_text = cell.text.strip()
#                 if cell_text not in cells:
#                     cells.append(cell_text)
#             # Ensure unique rows within the table
#             if cells not in rows:
#                 rows.append(cells)
#         return rows

#     # Function to add text and table content to a section
#     def add_content_to_section(section_name, content):
#         if section_name not in data["Facilitator_Guide"]:
#             data["Facilitator_Guide"][section_name] = []
#         # Check for duplication before adding content
#         if content not in data["Facilitator_Guide"][section_name]:
#             data["Facilitator_Guide"][section_name].append(content)

#     # Variables to track the current section
#     current_section = None

#     # Iterate through the elements of the document
#     for element in doc.element.body:
#         if isinstance(element, CT_P):  # It's a paragraph
#             para = Paragraph(element, doc)
#             text = para.text.strip()
#             if "Course Overview" in text or "Learning Outcomes" in text or "Lesson Plan" in text or "Structure & Duration" in text:
#                 # Use header content to set the current section
#                 current_section = text
#             elif text:
#                 add_content_to_section(current_section, text)
#         elif isinstance(element, CT_Tbl):  # It's a table
#             tbl = Table(element, doc)
#             table_content = parse_table(tbl)
#             if current_section:
#                 add_content_to_section(current_section, {"table": table_content})

#     return data

# def retrieveParsedJSON(raw_data, llm_config):
#     # 1. User Proxy Agent (Provides unstructured data to the interpreter)
#     user_proxy = UserProxyAgent(
#         name="User",
#         is_termination_msg=lambda msg: msg.get("content") is not None and "TERMINATE" in msg["content"],
#         human_input_mode="NEVER",  # Automatically provides unstructured data
#         code_execution_config={"work_dir": "output", "use_docker": False} # Takes data from a directory
#     )

#     # 2. Interpreter Agent (Converts unstructured data into structured data)
#     interpreter = AssistantAgent(
#         name="Interpreter",
#         llm_config=llm_config,
#         system_message="""
#         You are an AI assistant that extracts specific information from a JSON object containing a Course Proposal Form. Your task is to interpret the JSON data, regardless of its structure or any irrelevant information, and extract the required information accurately.

#         ---

#         **Task:** Extract the following information from the provided JSON data:

#         ### **Required Information:**

#         - **Course_Title**
#         - **TGS_Ref_No**
#         - **TSC_Title**
#         - **TSC_Code**
#         - **Abilities**: A list of ability statements (e.g., `["A1: ...", "A2: ..."]`)
#         - **Knowledge**: A list of knowledge statements (e.g., `["K1: ...", "K2: ..."]`)
#         - **Learning_Outcomes**: A list of learning outcomes (e.g., `["LO1: ...", "LO2: ..."]`)
#         - **Course_Outline**: A list of dictionaries, each representing a Learning Unit (LU) with its topics and subtopics.
#             - Each LU should include:
#                 - **Learning_Unit_Title**: Include the "LUx: " prefix.
#                 - **Topics**: A list of topics covered under the LU.
#                     - For each Topic:
#                         - **Topic_Title**: Include the "Topic x: " prefix and associated K and A statements in parentheses.
#                         - **Bullet_Points**: A list of bullet points under the topic.
#         - **Assessment_Methods**: A list of assessment methods, using full terms (e.g., `["Written Assessment - Short Answer Questions (WA-SAQ)", "Case Study (CS)"]`).

#         ---

#         **Instructions:**

#         - **Data Parsing**: Carefully parse the JSON data to locate and extract the required information.

#         - **Topic Structure Parsing**:

#         - **Learning Units (LUs)**:

#             - Identify each LU in the "Learning Outcomes" section under "table".
#             - The LU title starts with `"LUx: "`, where `x` is the LU number.

#         - **Topics within LUs**:

#             - **Identify Topics**:

#             - Look for lines starting with `"Topic x "`, where `x` is the topic number.
#             - The topic title follows `"Topic x "`, and any associated K and A statements are in parentheses.
#             - **Example**: `"Topic 1 Introduction to Copilot for Microsoft 365 (K1, A1)"`.

#             - **Bullet Points**:

#             - The lines following the topic title are bullet points for that topic.
#             - Collect these lines until you reach the next line that starts with:

#                 - `"Topic x "` (indicating the next topic),
#                 - `"LOx – "` (indicating the Learning Outcome), or
#                 - Another section header.

#             - **Validation**:

#             - Ensure all topics are extracted by checking the sequence of topic numbers.
#             - If topic numbers are skipped or duplicated, review the data carefully.

#         - **Learning Outcomes (LOs)**:

#             - The Learning Outcome starts with a line beginning with `"LOx – "`, where `x` is the LO number.

#         - **Knowledge (K) and Ability (A) Statements**:

#             - Following the LOs, K and A statements are listed, starting with `"Kx: "` and `"Ax: "`.

#         - **Irrelevant Data**: Ignore any irrelevant data or sections that are not needed for the extraction.

#         - **Consistency**: Ensure that all extracted information matches the data in the JSON input accurately.

#         - **Text Normalization**:

#             - Replace special characters:

#                 - En dashes (–) and em dashes (—) with regular hyphens (-).
#                 - Curly quotes (“ ”) with straight quotes (").
#                 - Non-ASCII characters with their closest ASCII equivalents.

#             - This normalization should be applied to all extracted text fields to ensure consistency.

#         - **Formatting**:

#             - **Time Fields**: Include units, e.g., "14 hrs", "2 hrs".
#             - **Learning Outcomes**: Include the "LOx: " prefix.
#             - **Knowledge and Abilities**: Include numbering (e.g., "K1: ...", "A1: ...").
#             - **Topic Titles**: Include the "Topic x: " prefix and associated K and A statements in parentheses.

#         - **Assessment Methods**:

#             - If only abbreviations are provided, replace them with their full terms:

#                 - "WA-SAQ" → "Written Assessment - Short Answer Questions (WA-SAQ)"
#                 - "CS" → "Case Study (CS)"
#                 - "PP" → "Practical Performance (PP)"
#                 - "OQ" → "Oral Questioning (OQ)"
#                 - "RP" → "Role Play (RP)"

#             - Use the standard assessment method abbreviations if needed.

#         - **Output Format**:

#             - Present the extracted information in a structured JSON format, using double quotes around all field names and string values.
#             - Remove any redundant or duplicate entries.
#             - Avoid any trailing commas.
#             - Ensure the JSON is well-formatted and valid.

#         - **Bullet Points Validation**:

#             - Ensure that all bullet points under each topic are fully extracted.
#             - If there is any discrepancy in the number of bullet points, re-extract them to match the source data.

#         - **Total Duration Validation**:

#             - Use the total duration specified in the data for "Total Course Duration Hours".
#             - Verify that this matches the sum of "Total Training Hours" and "Total Assessment Hours".
#             - If there is a discrepancy, use the total duration specified in the data as authoritative.

#         ---

#         **Expected Output Format**:

#         ```json
#         {
#             "Course_Title": "...",
#             "TGS_Ref_No": "...",
#             "TSC_Title": "...",
#             "TSC_Code": "...",
#             "Abilities": [
#                 "A1: ...",
#                 "A2: ...",
#                 "A3: ..."
#             ],
#             "Knowledge": [
#                 "K1: ...",
#                 "K2: ...",
#                 "K3: ..."
#             ],
#             "Learning_Outcomes": [
#                 "LO1: ...",
#                 "LO2: ...",
#                 "LO3: ..."
#             ],
#             "Course_Outline": [
#                 {
#                 "Learning_Unit_Title": "LU1: ...",
#                 "Topics": [
#                     {
#                     "Topic_Title": "Topic 1: ... (Kx, Ay)",
#                     "Bullet_Points": [
#                         "...",
#                         "..."
#                     ]
#                     }
#                     // Additional topics...
#                 ]
#                 }
#                 // Additional learning units...
#             ],
#             "Assessment_Methods": [
#                 "Written Assessment - Short Answer Questions (WA-SAQ)",
#                 "Case Study (CS)"
#             ]
#         }

#     """
#     )
#     chat_result = user_proxy.initiate_chat(
#         recipient=interpreter,
#         summary_method="last_msg",
#         max_turns=3,
#         message=f"""               
#         Please extract and structure the following data: {raw_data}.
#         **Return the extracted information as a complete JSON dictionary containing the specified fields. Do not truncate or omit any data. Include all fields and their full content. Do not use '...' or any placeholders to replace data.**
#         Simply return the JSON dictionary object directly and 'TERMINATE'.
#         """
#     )
 
#     response = chat_result.chat_history[-1].get("content", "")

#     return response

# def strip_response(response: str) -> str:
#     """
#     Strips triple backticks and `json` formatting hint from a response string.

#     Args:
#         response (str): The raw response encapsulated in triple backticks.

#     Returns:
#         str: The cleaned JSON string.
#     """
#     if response.startswith("```json") and response.endswith("```"):
#         # Remove triple backticks and the 'json' label
#         response = response.strip("```json").strip("```").strip()
#     # Deserialize JSON
#     try:
#         return json.loads(response)
#     except json.JSONDecodeError as e:
#         raise ValueError(f"Error parsing JSON response: {e}")
    
def parse_documents(doc_paths, output_dir="./parsed_documents"):
    """
    Parse multiple PDF documents into structured nodes for RAG workflows.
    
    Args:
        doc_paths (list of str): Paths to PDF documents to be parsed.
        output_dir (str): Directory to save parsed results.

    Returns:
        list: Parsed nodes for further processing.
    """
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Initialize LlamaParse
    llama_parser = LlamaParse(result_type="markdown", num_workers=8)
    
    # Process each document
    all_nodes = []
    for doc_path in doc_paths:
        print(f"Parsing document: {doc_path}")
        documents = llama_parser.load_data(doc_path)  # Parse the PDF
        
        # Convert documents into nodes
        node_parser = MarkdownElementNodeParser(llm=llm)
        nodes = node_parser.get_nodes_from_documents(documents)
        
        # Save parsed nodes to JSON
        output_file = Path(output_dir) / f"{Path(doc_path).stem}_parsed.json"
        with open(output_file, "w") as f:
            json.dump([node.to_dict() for node in nodes], f, indent=4)
        
        print(f"Saved parsed nodes to: {output_file}")
        all_nodes.extend(nodes)
    
    return all_nodes

def get_pdf_files(directory):
    pdf_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith('.pdf'):
                pdf_files.append(os.path.join(root, file))
    return pdf_files

input_fg = "./Slides/docs/FG_TGS-2024048313_Java Programming Methodologies_v2.docx"
input_directory = "./Slides/docs/REST/"

# parsed_data = parse_fg_document(input_fg)

# # Save to JSON file
# output_file = "parsed_fg.json"
# saveToJSON(output_file, parsed_data)

# response = retrieveParsedJSON(parsed_data, llm_config)
# context = strip_response(response)

with open('parsed_response.json', 'r') as f:
    context = json.load(f)

# Get all PDF files in the directory
pdf_files = get_pdf_files(input_directory)

# Parse the documents and save results
nodes = parse_documents(pdf_files)
print(f"Parsed {len(nodes)} nodes from {len(pdf_files)} documents.")
recursive_index = VectorStoreIndex(nodes=nodes)

# Reranker
rerank = FlagEmbeddingReranker(model="BAAI/bge-reranker-large", top_n=5)

# Configure Query Engine
recursive_query_engine = recursive_index.as_query_engine(
    similarity_top_k=5, node_postprocessors=[rerank], verbose=True
)

# Define the updated query format
query_template = """Retrieve key information on the following topic from the vector store ONLY. Ensure all responses are grounded strictly in retrieved data.

### Instructions
1. Retrieve information that specifically corresponds to {subtopic} with context of the broarder topic {topic}, matching the relevant chapters, sections, or locations in the vector store.
2. If no relevant information exists for the subtopic in the vector store, return the response "NO CONTEXT".
3. Ensure that each response provides content grounded only in the vector store, without additional extrapolation.
3. Organize information concisely, in bullet points.

**Output Format:**
{{
    "subtopic": "{subtopic}",
    "keypoints": [
        {{
            "keypoint": "Key Point 1 [Specific section of the subtopic]",
            "bullets": [
                "Bullet point 1: [Specific lecture content on Key Point 1]",
                "Bullet point 2: [Another specific lecture content on Key Point 1]"
            ]
        }},
        {{
            "keypoint": "Key Point 2",
            "bullets": [
                "Bullet point 1: [Specific lecture content about Key Point 2]",
                "Bullet point 2: [Another specific lecture content about Key Point 2]"
            ]
        }}
    ]
}}

**Requirements:**
- Ensure each bullet point provides relevant details **about the content of the keypoint** itself, not references to a source.
- Provide no more than two sentences per bullet point, ensuring concise, detailed information directly on the keypoint topic, not whether it is covered in Chapter X.
"""

context["Topics_Outline"] = []  # Initialize results as a dictionary with a "topics" key

# Run the query for each Learning Unit and Topic
for lu in context["Course_Outline"]:
    # Create a learning unit dictionary to hold topics and subtopics
    lu_data = {
        "learning_unit": lu["Learning_Unit_Title"],
        "topics": []
    }

    for topic in lu["Topics"]:
        # Format the query with topic-specific information
        topic_title = topic["Topic_Title"]
        bullet_points = topic["Bullet_Points"]

        # Add the topic and subtopic data to the learning unit
        topic_data = {
            "topic_title": topic_title,
            "keypoints": []
        }

        # Query for each bullet point as a subtopic
        for subtopic in bullet_points:
            query = query_template.format(subtopic=subtopic, topic=topic_title)
            
            # Execute the query
            response = recursive_query_engine.query(query)
            
            # Extract and process JSON response
            try:
                json_data = eval(response.response)  # Convert response string to dictionary
                
                # Add subtopic keypoints to the topic data
                subtopic_data = {
                    "subtopic": json_data.get("subtopic"),
                    "keypoints": json_data.get("keypoints", [])
                }
                topic_data["keypoints"].append(subtopic_data)

            except (SyntaxError, ValueError):
                print(f"Failed to parse response for subtopic: {subtopic}")
            
            print(f"\n***********Processing {topic_title} - Subtopic: {subtopic}***********")

        # Append the completed topic data to the learning unit's topics
        lu_data["topics"].append(topic_data)

    # Append the completed learning unit data to the topics outline
    context["Topics_Outline"].append(lu_data)
    
saveToJSON("rag_done_output_v2.json", context)
print(f"\n***********Generating Slide Deck for {context['Course_Title']}***********")

# with open('rag_done_output.json', 'r') as f:
#     context = json.load(f)

# Create the content generator agent
slides_writer = AssistantAgent(
    name="Slides_Generator",
    system_message=f"""
    You are responsible for generating PowerPoint slides dynamically based on the provided context.
    The `context` JSON will be included in the messages you receive.
    
    Instructions:
    - Extract the `context` JSON from the incoming message.
    - Pass the `context` as a parameter to the `generate_slides` function.
    - Confirm completion by returning `TERMINATE`.
    """,
    llm_config=llm_config,
)

# Create the user proxy agent
user_proxy = UserProxyAgent(
    name="User_Proxy",
    human_input_mode="NEVER",
    max_consecutive_auto_reply=2,
    # is_termination_msg=lambda x: x.get("content", "").find("TERMINATE") >= 0,
    code_execution_config={"work_dir": "output", "use_docker": False} 
)

@user_proxy.register_for_execution()
@slides_writer.register_for_llm(name="generate_slides", description="Generate the Master Slides")
def generate_slides(context: dict) -> str:
    """
    Dynamically generate slides based on the passed context.
    """
    # Function to build replacements dictionary from JSON data
    def build_replacements(data):
        replacements = {}
        for key, value in data.items():
            placeholder = f"<{key.upper()}>"
            replacements[placeholder] = value
        return replacements
    
    # Function to replace placeholders
    def replace_placeholder(text, replacements):
        for placeholder, replacement in replacements.items():
            if isinstance(replacement, dict):
                # Skip dictionaries
                continue
            elif isinstance(replacement, list):
                # Check if all items are strings
                if all(isinstance(item, str) for item in replacement):
                    replacement_text = "\n".join(replacement)
                    text = text.replace(placeholder, replacement_text)
                else:
                    # Skip lists that contain non-string items (e.g., lists of dicts)
                    continue
            else:
                text = text.replace(placeholder, str(replacement))
        return text

    # Function to add bullet points from a list to a text frame
    def add_bullet_points(text_frame, bullet_points):
        text_frame.clear()  # Clear any existing text
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Level 0 for top-level bullet points

    def get_slide_number_placeholder_properties():
        # Positioning and formatting based on the data provided in the images
        position = (Cm(23.54), Cm(12.95), Cm(1.52), Cm(1.09))  # Left, Top, Width, Height in centimeters
        font_properties = {
            'font_size': Pt(10),
            'font_name': "Arial",
            'font_color': RGBColor(89, 89, 89),  # Dark Gray, Background 2
            'alignment': 'middle',  # Middle alignment
            'horizontal_alignment': 'right',  # Right horizontal alignment
            'autofit': False,
            'wrap_text': True,
            'margins': (Cm(0.25), Cm(0.25), Cm(0.25), Cm(0.25))  # Left, Right, Top, Bottom margins in centimeters
        }
        return position, font_properties

    def move_slides(prs1, prs2, placeholder_text):
        """Move slides from prs2 to prs1 after the placeholder slide identified by placeholder_text."""
        placeholder_index = None
        for idx, slide in enumerate(prs1.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and placeholder_text in shape.text:
                    placeholder_index = idx
                    break
            if placeholder_index is not None:
                break

        if placeholder_index is None:
            print(f"Placeholder '{placeholder_text}' not found.")
            return prs1, []

        # Collect slides after the placeholder
        slides_after_placeholder = [prs1.slides[i] for i in range(placeholder_index + 1, len(prs1.slides))]

        # Remove slides after the placeholder
        for i in range(len(prs1.slides) - 1, placeholder_index, -1):
            rId = prs1.slides._sldIdLst[i].rId
            prs1.part.drop_rel(rId)
            del prs1.slides._sldIdLst[i]

        # Insert slides from prs2 after the placeholder
        subtopic_slides = []
        for slide in prs2.slides:
            new_slide = prs1.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                sp = shape.element
                sp.getparent().remove(sp)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_stream = BytesIO(shape.image.blob)
                    new_slide.shapes.add_picture(
                        image_stream,
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                else:
                    new_shape = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

            try:
                new_slide.shapes.title.text = slide.shapes.title.text
            except AttributeError:
                pass

            subtopic_slides.append(new_slide)

        # Re-add slides that were after the placeholder
        for slide in slides_after_placeholder:
            new_slide = prs1.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                sp = shape.element
                sp.getparent().remove(sp)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_stream = BytesIO(shape.image.blob)
                    new_slide.shapes.add_picture(
                        image_stream,
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                else:
                    new_shape = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

            try:
                new_slide.shapes.title.text = slide.shapes.title.text
            except AttributeError:
                pass

        # Delete the placeholder slide
        rId = prs1.slides._sldIdLst[placeholder_index].rId
        prs1.part.drop_rel(rId)
        del prs1.slides._sldIdLst[placeholder_index]
        
        return prs1, subtopic_slides

    def generate_course_outline(prs2, context, content_slide_layout):
        course_outline = context["<COURSE_OUTLINE>"]
        topics_per_slide = 2
        learning_units = course_outline

        topic_entries = []
        for lu in learning_units:
            lu_title = lu['Learning_Unit_Title']
            for topic in lu['Topics']:
                topic_title = topic['Topic_Title']
                bullet_points = topic['Bullet_Points']
                topic_entries.append((topic_title, bullet_points))

        for i in range(0, len(topic_entries), topics_per_slide):
            slide = prs2.slides.add_slide(content_slide_layout)
            title_placeholder = slide.shapes.title
            title_placeholder.text = "Course Outline"

            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for entry in topic_entries[i:i + topics_per_slide]:
                topic_title, bullet_points = entry
                p = text_frame.add_paragraph()
                p.text = f"{topic_title}"
                p.font.bold = True
                for bullet in bullet_points:
                    sub_p = text_frame.add_paragraph()
                    sub_p.text = bullet
                    sub_p.level = 1

    # Function to generate topic slides
    def generate_topics(prs2, context, content_slide_layout, title_slide_layout):
        topics_outline = context['<TOPICS_OUTLINE>']
        for learning_unit in topics_outline:
            lu_title = learning_unit['learning_unit']
            for topic in learning_unit['topics']:
                topic_title = topic['topic_title']
                # Add a topic slide
                topic_slide = prs2.slides.add_slide(title_slide_layout)
                topic_title_placeholder = topic_slide.shapes.title
                topic_title_placeholder.text = f"{topic_title}"

                for subtopic in topic['keypoints']:
                    subtopic_title = subtopic['subtopic']
                    # Add a subtopic slide
                    slide = prs2.slides.add_slide(content_slide_layout)
                    title_placeholder = slide.shapes.title
                    title_placeholder.text = subtopic_title

                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()

                    for keypoint in subtopic['keypoints']:
                        keypoint_text = keypoint['keypoint']
                        bullets = keypoint['bullets']

                        if keypoint_text == "NO CONTEXT":
                            continue

                        p = text_frame.add_paragraph()
                        p.text = keypoint_text
                        p.font.bold = True
                        for bullet in bullets:
                            if bullet != "NO CONTEXT":
                                bullet_p = text_frame.add_paragraph()
                                bullet_p.text = bullet

    # Function to update slide numbers
    def update_slide_numbers(presentation, subtopic_slides, slide_number_position, font_properties):
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_number_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    slide_number_placeholder = shape
                    break
            if slide_number_placeholder:
                slide_number_placeholder.text = str(idx)
            elif slide in subtopic_slides:
                left, top, width, height = slide_number_position
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = str(idx)
                font = p.font
                font.size = font_properties['font_size']
                font.name = font_properties['font_name']
                font.color.rgb = font_properties['font_color']
                tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = font_properties['margins']
                tf.word_wrap = font_properties['wrap_text']
                
                # Set auto_size based on font_properties['autofit']
                tf.auto_size = MSO_AUTO_SIZE.NONE if not font_properties['autofit'] else MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                # Set vertical alignment to middle
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Set paragraph alignment to right
                p.alignment = PP_ALIGN.RIGHT

    def remove_unused_placeholders(presentation):
        for slide in presentation.slides:
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.is_placeholder:
                    # Check if the placeholder is empty
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if not text:
                            shapes_to_remove.append(shape)
                    elif shape.placeholder_format.type == MSO_SHAPE_TYPE.PICTURE:
                        # For picture placeholders, check if an image is inserted
                        if not shape.has_image:
                            shapes_to_remove.append(shape)
                    else:
                        # Other placeholder types can be checked similarly
                        shapes_to_remove.append(shape)
            # Remove the unused placeholders
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)

    # Main function logic
    try:
        # Load the PowerPoint template
        template_path = 'Slides/templates/Non-WSQ/(Template) SG - Master Trainer Slides - Course Title - version.pptx'
        prs1 = Presentation(template_path)
        content_slide_layout = prs1.slide_layouts[1]
        title_slide_layout = prs1.slide_masters[0].slide_layouts.get_by_name("SECTION_HEADER")
        replacements = build_replacements(context)

        # Step 1: Replace placeholders on existing slides
        for slide in prs1.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "<LEARNING_OUTCOMES>" in shape.text:
                        add_bullet_points(shape.text_frame, replacements["<LEARNING_OUTCOMES>"])
                    elif "<ASSESSMENT_METHODS>" in shape.text:
                        add_bullet_points(shape.text_frame, replacements["<ASSESSMENT_METHODS>"])
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.text = replace_placeholder(paragraph.text, replacements)

        # Step 2: Generate Course Outline Slides
        prs2 = Presentation()
        generate_course_outline(prs2, replacements, content_slide_layout)

        # Insert the course outline slides after the <COURSE_OUTLINE> placeholder
        prs1, outline_slides = move_slides(prs1, prs2, "<COURSE_OUTLINE>")
        # Update slide numbers
        slide_number_position, font_properties = get_slide_number_placeholder_properties()
        update_slide_numbers(prs1, outline_slides, slide_number_position, font_properties)


        # Step 3: Generate Topic Slides
        prs2 = Presentation()  # Reset prs2 for topics
        generate_topics(prs2, replacements, content_slide_layout, title_slide_layout)

        # Insert the topic slides after the <TOPICS_OUTLINE> placeholder
        prs1, subtopic_slides = move_slides(prs1, prs2, "<TOPICS_OUTLINE>")

        # Update slide numbers
        slide_number_position, font_properties = get_slide_number_placeholder_properties()
        update_slide_numbers(prs1, subtopic_slides, slide_number_position, font_properties)

        # Step 4: Remove unused placeholders
        remove_unused_placeholders(prs1)

        # Save final presentation
        output_path = 'Slides/output/autogen_output.pptx'
        prs1.save(output_path)
        print(f"Slides successfully generated at {output_path}")
        return output_path

    except Exception as e:
        print(f"Error generating slides: {e}")
        return str(e)

user_proxy.initiate_chat(
    slides_writer,
    message=f"""
    Please generate the PowerPoint slides using the following context:
    {json.dumps(context)}
    """
)
